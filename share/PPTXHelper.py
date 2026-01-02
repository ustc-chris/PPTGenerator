from pptx import Presentation

# --- Namespaces used by PowerPoint sections (p14) ---
_NS = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "p14": "http://schemas.microsoft.com/office/powerpoint/2010/main",
}

def _remove_slide_from_sections(prs, slide_id: int, remove_empty_sections: bool = True):
    pres = prs.part._element

    # 找到所有 <sldId id="slide_id">
    sldid_elems = pres.xpath(
        f".//*[local-name()='sectionLst']"
        f"//*[local-name()='section']"
        f"//*[local-name()='sldIdLst']"
        f"/*[local-name()='sldId' and @id='{slide_id}']"
    )

    for sldid in sldid_elems:
        sldid.getparent().remove(sldid)

    if not remove_empty_sections:
        return

    # 删除空的 section
    empty_sections = pres.xpath(
        ".//*[local-name()='section' and "
        "count(.//*[local-name()='sldId'])=0]"
    )
    for sec in empty_sections:
        sec.getparent().remove(sec)

    # 删除空的 sectionLst
    empty_sectionlsts = pres.xpath(
        ".//*[local-name()='sectionLst' and "
        "count(.//*[local-name()='section'])=0]"
    )
    for lst in empty_sectionlsts:
        lst.getparent().remove(lst)


def delete_slide_by_index_with_sections(prs, n: int):
    """
    n = 0, 1, ..., len(prs.slides) - 1
    """
    slides = prs.slides
    total = len(slides)
    if n < 0:
        n += total
    if n < 0 or n >= total:
        raise IndexError(n)

    slide = slides[n]
    slide_id = slide.slide_id

    # 1) 先修 Sections
    _remove_slide_from_sections(prs, slide_id)

    # 2) 再删 slide 本体
    sldIdLst = slides._sldIdLst
    target = None
    for sldId in sldIdLst:
        if sldId.id == slide_id:
            target = sldId
            break

    if target is None:
        raise RuntimeError("slide not found")

    prs.part.drop_rel(target.rId)
    sldIdLst.remove(target)

def clear_all_slides_with_sections(prs):
    # 先清空所有 sections（最干净）
    pres = prs.part._element
    for lst in pres.xpath(".//*[local-name()='sectionLst']"):
        lst.getparent().remove(lst)

    # 再从后往前删 slides
    for i in range(len(prs.slides) - 1, -1, -1):
        slide = prs.slides[i]
        slide_id = slide.slide_id

        sldIdLst = prs.slides._sldIdLst
        target = None
        for sldId in sldIdLst:
            if sldId.id == slide_id:
                target = sldId
                break
        if target is None:
            continue

        prs.part.drop_rel(target.rId)
        sldIdLst.remove(target)


from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

P14_NS = "http://schemas.microsoft.com/office/powerpoint/2010/main"

def _get_or_create_sectionlst(prs):
    pres = prs.part._element

    # 已有 sectionLst 就直接返回（不管是 p14 还是别的前缀）
    lsts = pres.xpath(".//*[local-name()='sectionLst']")
    if lsts:
        return lsts[0]

    # 确保 extLst 存在
    extlst = pres.xpath(".//*[local-name()='extLst']")
    if extlst:
        extlst = extlst[0]
    else:
        extlst = parse_xml(f"<p:extLst {nsdecls('p')}/>")
        pres.append(extlst)

    # 新建 ext + sectionLst（显式声明 xmlns:p14，避免 nsdecls('p14')）
    ext = parse_xml(
        f"""
        <p:ext {nsdecls('p')} xmlns:p14="{P14_NS}"
            uri="{{EFAFB233-063F-42B5-8137-9DF3F51BA10A}}">
          <p14:sectionLst/>
        </p:ext>
        """
    )
    extlst.append(ext)
    return ext.xpath(".//*[local-name()='sectionLst']")[0]


def add_section(prs, name: str, slide_indices):
    """
    给 PPT 添加一个分节
    :param name: 分节名（PowerPoint 左侧显示）
    :param slide_indices: slide 的索引列表（0-based）
    """
    slides = prs.slides
    sectionlst = _get_or_create_sectionlst(prs)

    # 收集 slide_id
    slide_ids = [slides[i].slide_id for i in slide_indices]

    # section id：PowerPoint 用一个整数，唯一即可
    used_ids = [
        int(sec.get("id"))
        for sec in sectionlst.xpath(".//*[local-name()='section']")
        if sec.get("id") is not None
    ]
    sec_id = max(used_ids, default=255) + 1

    # 构造 section XML
    sld_ids_xml = "\n".join(
        f'<p14:sldId id="{sid}"/>' for sid in slide_ids
    )

    section = parse_xml(
        f"""
        <p14:section xmlns:p14="{P14_NS}"
            name="{name}" id="{sec_id}">
          <p14:sldIdLst>
            {sld_ids_xml}
          </p14:sldIdLst>
        </p14:section>
        """
    )

    sectionlst.append(section)