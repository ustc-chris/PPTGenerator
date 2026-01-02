import os
import pptx 
from pptx import Presentation
from enum import Enum
import json
from .Template import Manager, TemplateType
from .PPTXHelper import delete_slide_by_index_with_sections
from .PPTXHelper import clear_all_slides_with_sections
from .PPTXHelper import add_section


class Lyric:
    def __init__(self):
        self.template = TemplateType.Undefined
        self.text = ""
        self.page = 0

class PPT:
    def __init__(self, template_ppt_file):
        self.lyrics = []
        self.ppt = Presentation(template_ppt_file)
        self.sections = {}

    def read(self, json_file=None, input_map=None):
        manager = Manager()
        if json_file:
            with open(json_file, 'r') as file:
                data = json.load(file)
        if input_map:
            data = input_map

        # 遍历
        page_index = 1
        for slide in data['slides']:
            for item in slide['items']:
                lyric = Lyric()
                lyric.template = slide['template']
                lyric.page = page_index
                manager.parse(lyric.template, item, lyric)
                self.lyrics.append(lyric)
                page_index += 1
    
    def add_section(self, page_index, section_name):
        self.sections[page_index] = section_name

    def write(self, filename):
        manager = Manager()
        for lyric in self.lyrics:
            manager.print(self.ppt, lyric)
        
        section_pages, section_names = [], []
        for page_index, section_name in sorted(self.sections.items()):
            section_pages.append(page_index)
            section_names.append(section_name)

        for i in range(len(section_pages)):
            if i == len(section_pages) - 1:  # 最后一个了
                pages = list(range(section_pages[i] - 1, len(self.ppt.slides))) # page - 1, page, ..., N-1
            else:
                pages = list(range(section_pages[i] - 1, section_pages[i+1] - 1)) # page - 1, page, ..., page+1 - 1 - 1
            add_section(self.ppt, section_names[i], pages)

        self.ppt.save(filename)

    def set_section(self, page_index, section_name):
        self.sections[page_index] = section_name

    def reset(self):
        clear_all_slides_with_sections(self.ppt)
        self.sections = {}


class PPTReader:
    def __init__(self, ppt_file):
        self.ppt = Presentation(ppt_file)
        self.NPages = len(self.ppt.slides)
        self.lyrics = []
        all_layouts = [layout.name for layout in self.ppt.slide_layouts]
        manager = Manager()
        for slide in self.ppt.slides:
            type_template = slide.slide_layout.name
            template_index = all_layouts.index(type_template)
            lyric = manager.convert_back(slide, template_index)
            self.lyrics.append(lyric)
        
    def convert_back(self):
        return {"slides": self.lyrics}
