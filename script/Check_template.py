from pptx import Presentation

# prs = Presentation("./share/template.pptx")
prs = Presentation("./梁祝侧屏PPTv2（内容修改+字体修改+分节+字体嵌入）/梁祝侧屏PPTv2（内容修改+字体修改+分节+字体嵌入）.pptx")

for i, layout in enumerate(prs.slide_layouts):
    print("=" * 60)
    print(f"Layout {i}: {layout.name}")

    for ph in layout.placeholders:
        print(f"  idx={ph.placeholder_format.idx:<2} | "
              f"type={ph.placeholder_format.type} | "
              f"name={ph.name}")